"""
מודול לטיפול בסוגי קבצים שונים והמרתם לטקסט
"""
import os
import logging
from typing import Dict, Any, Optional, Tuple
import mimetypes
import traceback

# ספריות לטיפול בסוגי קבצים שונים
from pypdf import PdfReader
import docx
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import html2text

# הגדרת לוגר
logger = logging.getLogger(__name__)

class FileParser:
    """
    מחלקה לטיפול בסוגי קבצים שונים והמרתם לטקסט
    """
    
    @staticmethod
    def parse_file(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        פונקציה ראשית לפרסור קובץ - מזהה את סוג הקובץ ומפעילה את הפרסר המתאים
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            טקסט מהקובץ ומטא-דאטה
        """
        try:
            # זיהוי סוג הקובץ לפי הסיומת
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # מטא-דאטה בסיסי
            metadata = {
                'filename': os.path.basename(file_path),
                'file_extension': file_ext,
                'file_size_bytes': os.path.getsize(file_path),
                'mime_type': mimetypes.guess_type(file_path)[0] or 'application/octet-stream'
            }
            
            # בחירת פרסר מתאים לפי סוג הקובץ
            if file_ext == '.pdf':
                return FileParser.parse_pdf(file_path, metadata)
            elif file_ext == '.docx':
                return FileParser.parse_docx(file_path, metadata)
            elif file_ext == '.xlsx':
                return FileParser.parse_excel(file_path, metadata)
            elif file_ext == '.pptx':
                return FileParser.parse_pptx(file_path, metadata)
            elif file_ext in ['.html', '.htm']:
                return FileParser.parse_html(file_path, metadata)
            elif file_ext in ['.txt', '.md', '.py', '.js', '.css', '.json', '.xml', '.csv']:
                # קבצי טקסט פשוטים
                return FileParser.parse_text(file_path, metadata)
            else:
                # ניסיון לפרסר כקובץ טקסט
                try:
                    return FileParser.parse_text(file_path, metadata)
                except Exception:
                    raise ValueError(f"סוג קובץ לא נתמך: {file_ext}")
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ {file_path}: {str(e)}")
            logger.error(traceback.format_exc())
            raise
    
    @staticmethod
    def parse_pdf(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ PDF"""
        try:
            reader = PdfReader(file_path)
            text = ""
            
            # מטא-דאטה נוסף
            metadata['page_count'] = len(reader.pages)
            if reader.metadata:
                for key, value in reader.metadata.items():
                    if key.startswith('/'):
                        clean_key = key[1:]  # הסרת ה-/ מהתחלת המפתח
                        metadata[f'pdf_{clean_key}'] = value
            
            # חילוץ טקסט מכל עמוד
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text:
                    text += f"\n--- עמוד {i+1} ---\n{page_text}\n"
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ PDF {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def parse_docx(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ Word (DOCX)"""
        try:
            doc = docx.Document(file_path)
            
            # מטא-דאטה נוסף
            metadata['paragraph_count'] = len(doc.paragraphs)
            core_props = doc.core_properties
            if core_props:
                metadata['docx_title'] = core_props.title
                metadata['docx_author'] = core_props.author
                metadata['docx_created'] = str(core_props.created) if core_props.created else None
                metadata['docx_modified'] = str(core_props.modified) if core_props.modified else None
            
            # חילוץ טקסט מכל פסקה
            text = "\n".join([para.text for para in doc.paragraphs if para.text])
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ Word {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def parse_excel(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ Excel (XLSX)"""
        workbook = None
        try:
            # ניסיון לפתוח את הקובץ עם טיפול בשגיאת קובץ נעול
            try:
                workbook = load_workbook(file_path, read_only=True, data_only=True)
            except PermissionError as pe:
                logger.error(f"הקובץ {file_path} נעול על ידי תהליך אחר: {str(pe)}")
                raise ValueError(f"הקובץ נעול על ידי תהליך אחר. אנא סגור את הקובץ אם הוא פתוח בתוכנה אחרת ונסה שוב.")
            
            # מטא-דאטה נוסף
            metadata['sheet_count'] = len(workbook.sheetnames)
            metadata['sheet_names'] = workbook.sheetnames
            
            # חילוץ טקסט מכל גיליון
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- גיליון: {sheet_name} ---\n"
                
                # הוספת תוכן התאים
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_values):  # רק שורות שיש בהן תוכן
                        text += " | ".join(row_values) + "\n"
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ Excel {file_path}: {str(e)}")
            raise
        finally:
            # סגירת ה-workbook אם הוא נפתח בהצלחה
            if workbook:
                try:
                    workbook.close()
                except Exception as close_error:
                    logger.warning(f"שגיאה בסגירת קובץ Excel {file_path}: {str(close_error)}")
    
    @staticmethod
    def parse_pptx(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ PowerPoint (PPTX)"""
        try:
            presentation = Presentation(file_path)
            
            # מטא-דאטה נוסף
            metadata['slide_count'] = len(presentation.slides)
            
            # חילוץ טקסט מכל שקופית
            text = ""
            for i, slide in enumerate(presentation.slides):
                text += f"\n--- שקופית {i+1} ---\n"
                
                # חילוץ טקסט מכל תיבת טקסט בשקופית
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                text += "\n".join(slide_text) + "\n"
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ PowerPoint {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def parse_html(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ HTML"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # שימוש ב-BeautifulSoup לפרסור ה-HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # מטא-דאטה נוסף
            title_tag = soup.find('title')
            if title_tag:
                metadata['html_title'] = title_tag.text
            
            # המרה לטקסט פשוט
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            text = h.handle(html_content)
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ HTML {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def parse_text(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ טקסט פשוט"""
        try:
            # ניסיון לפתוח עם קידוד UTF-8
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                # אם נכשל, ננסה עם קידוד אחר
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()
            
            # מטא-דאטה נוסף
            metadata['line_count'] = text.count('\n') + 1
            metadata['char_count'] = len(text)
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ טקסט {file_path}: {str(e)}")
            raise 