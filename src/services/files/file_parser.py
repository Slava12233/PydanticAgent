"""
מודול לטיפול בסוגי קבצים שונים והמרתם לטקסט
"""

import os
import logging
from typing import Dict, Any, Optional, Tuple
import mimetypes
import traceback
from pathlib import Path
from datetime import datetime

# ספריות לטיפול בסוגי קבצים שונים
from pypdf import PdfReader
import docx
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import html2text

from src.utils.logger import setup_logger

logger = setup_logger(__name__)

class FileParser:
    """מחלקה לטיפול בסוגי קבצים שונים והמרתם לטקסט"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'parse_pdf',
        '.docx': 'parse_docx',
        '.xlsx': 'parse_excel',
        '.pptx': 'parse_pptx',
        '.html': 'parse_html',
        '.htm': 'parse_html',
        '.txt': 'parse_text',
        '.md': 'parse_text',
        '.py': 'parse_text',
        '.js': 'parse_text',
        '.css': 'parse_text',
        '.json': 'parse_text',
        '.xml': 'parse_text',
        '.csv': 'parse_text'
    }
    
    @classmethod
    def parse_file(cls, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """פונקציה ראשית לפרסור קובץ - מזהה את סוג הקובץ ומפעילה את הפרסר המתאים
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            טקסט מהקובץ ומטא-דאטה
            
        Raises:
            ValueError: אם סוג הקובץ לא נתמך
            FileNotFoundError: אם הקובץ לא קיים
            PermissionError: אם אין הרשאות לקובץ
        """
        try:
            # בדיקת קיום הקובץ
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"הקובץ {file_path} לא קיים")
                
            # בדיקת הרשאות
            if not os.access(file_path, os.R_OK):
                raise PermissionError(f"אין הרשאות קריאה לקובץ {file_path}")
            
            # זיהוי סוג הקובץ לפי הסיומת
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # מטא-דאטה בסיסי
            file_stat = os.stat(file_path)
            metadata = {
                'filename': os.path.basename(file_path),
                'file_extension': file_ext,
                'file_size_bytes': file_stat.st_size,
                'mime_type': mimetypes.guess_type(file_path)[0] or 'application/octet-stream',
                'created_at': datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
                'modified_at': datetime.fromtimestamp(file_stat.st_mtime).isoformat()
            }
            
            # בדיקה אם הסוג נתמך
            if file_ext not in cls.SUPPORTED_EXTENSIONS:
                raise ValueError(f"סוג הקובץ {file_ext} לא נתמך")
                
            # הפעלת הפרסר המתאים
            parser_method = getattr(cls, cls.SUPPORTED_EXTENSIONS[file_ext])
            return parser_method(file_path, metadata)
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור הקובץ {file_path}: {str(e)}")
            logger.error(traceback.format_exc())
            raise
            
    @staticmethod
    def parse_pdf(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ PDF
        
        Args:
            file_path: נתיב לקובץ
            metadata: מטא-דאטה בסיסי
            
        Returns:
            טקסט מהקובץ ומטא-דאטה מורחב
        """
        try:
            reader = PdfReader(file_path)
            text = ""
            
            # מטא-דאטה נוסף
            metadata['page_count'] = len(reader.pages)
            if reader.metadata:
                for key, value in reader.metadata.items():
                    if key.startswith('/'):
                        clean_key = key[1:]  # הסרת ה-/ מהתחלת המפתח
                        metadata[f'pdf_{clean_key}'] = value
                        
            # חילוץ טקסט מכל עמוד
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text:
                    text += f"\n--- עמוד {i+1} ---\n{page_text}\n"
                    
            return text, metadata
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ PDF {file_path}: {str(e)}")
            raise
            
    @staticmethod
    def parse_docx(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ Word (DOCX)
        
        Args:
            file_path: נתיב לקובץ
            metadata: מטא-דאטה בסיסי
            
        Returns:
            טקסט מהקובץ ומטא-דאטה מורחב
        """
        try:
            doc = docx.Document(file_path)
            
            # מטא-דאטה נוסף
            metadata['paragraph_count'] = len(doc.paragraphs)
            core_props = doc.core_properties
            if core_props:
                metadata['docx_title'] = core_props.title
                metadata['docx_author'] = core_props.author
                metadata['docx_created'] = str(core_props.created) if core_props.created else None
                metadata['docx_modified'] = str(core_props.modified) if core_props.modified else None
                
            # חילוץ טקסט מכל פסקה
            text = "\n".join([para.text for para in doc.paragraphs if para.text])
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ Word {file_path}: {str(e)}")
            raise
            
    @staticmethod
    def parse_excel(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ Excel (XLSX)
        
        Args:
            file_path: נתיב לקובץ
            metadata: מטא-דאטה בסיסי
            
        Returns:
            טקסט מהקובץ ומטא-דאטה מורחב
        """
        workbook = None
        try:
            # ניסיון לפתוח את הקובץ עם טיפול בשגיאת קובץ נעול
            try:
                workbook = load_workbook(file_path, read_only=True, data_only=True)
            except PermissionError as pe:
                logger.error(f"הקובץ {file_path} נעול על ידי תהליך אחר: {str(pe)}")
                raise ValueError("הקובץ נעול על ידי תהליך אחר. אנא סגור את הקובץ אם הוא פתוח בתוכנה אחרת ונסה שוב.")
                
            # מטא-דאטה נוסף
            metadata['sheet_count'] = len(workbook.sheetnames)
            metadata['sheet_names'] = workbook.sheetnames
            
            # חילוץ טקסט מכל גיליון
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- גיליון: {sheet_name} ---\n"
                
                # הוספת תוכן התאים
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_values):  # רק שורות שיש בהן תוכן
                        text += " | ".join(row_values) + "\n"
                        
            return text, metadata
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ Excel {file_path}: {str(e)}")
            raise
        finally:
            # סגירת ה-workbook אם הוא נפתח בהצלחה
            if workbook:
                try:
                    workbook.close()
                except Exception as close_error:
                    logger.warning(f"שגיאה בסגירת קובץ Excel {file_path}: {str(close_error)}")
                    
    @staticmethod
    def parse_pptx(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ PowerPoint (PPTX)
        
        Args:
            file_path: נתיב לקובץ
            metadata: מטא-דאטה בסיסי
            
        Returns:
            טקסט מהקובץ ומטא-דאטה מורחב
        """
        try:
            presentation = Presentation(file_path)
            
            # מטא-דאטה נוסף
            metadata['slide_count'] = len(presentation.slides)
            
            # חילוץ טקסט מכל שקופית
            text = ""
            for i, slide in enumerate(presentation.slides):
                text += f"\n--- שקופית {i+1} ---\n"
                
                # חילוץ טקסט מכל תיבת טקסט בשקופית
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                        
                text += "\n".join(slide_text) + "\n"
                
            return text, metadata
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ PowerPoint {file_path}: {str(e)}")
            raise
            
    @staticmethod
    def parse_html(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ HTML
        
        Args:
            file_path: נתיב לקובץ
            metadata: מטא-דאטה בסיסי
            
        Returns:
            טקסט מהקובץ ומטא-דאטה מורחב
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
                
            # שימוש ב-BeautifulSoup לפרסור ה-HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # מטא-דאטה נוסף
            title_tag = soup.find('title')
            if title_tag:
                metadata['html_title'] = title_tag.text
                
            # המרה לטקסט פשוט
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            text = h.handle(html_content)
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ HTML {file_path}: {str(e)}")
            raise
            
    @staticmethod
    def parse_text(file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """פרסור קובץ טקסט פשוט
        
        Args:
            file_path: נתיב לקובץ
            metadata: מטא-דאטה בסיסי
            
        Returns:
            טקסט מהקובץ ומטא-דאטה מורחב
        """
        try:
            # ניסיון לפתוח עם קידוד UTF-8
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                # אם נכשל, ננסה עם קידוד אחר
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()
                    
            # מטא-דאטה נוסף
            metadata['line_count'] = text.count('\n') + 1
            metadata['char_count'] = len(text)
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"שגיאה בפרסור קובץ טקסט {file_path}: {str(e)}")
            raise 